import streamlit as st
import pdfplumber
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from openai import OpenAI
import json
from typing import List, Dict
import os
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.docstore.document import Document
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import base64
import seaborn as sns
import matplotlib.pyplot as plt

class BusinessAnalyzer:
    def __init__(self, openai_client):
        self.client = openai_client

    def generate_insights(self, data: pd.DataFrame) -> Dict:
        # Convert DataFrame to string representation
        data_str = data.to_string()
        
        prompt = f"""
        Analyze the following medical document data and provide comprehensive business insights:
        {data_str}

        Please provide insights in the following JSON format:
        {{
            "key_findings": [
                "finding 1",
                "finding 2"
            ],
            "business_impact": [
                "impact 1",
                "impact 2"
            ],
            "recommendations": [
                "recommendation 1",
                "recommendation 2"
            ],
            "risk_factors": [
                "risk 1",
                "risk 2"
            ],
            "opportunities": [
                "opportunity 1",
                "opportunity 2"
            ]
        }}
        """
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-4-turbo-preview",
                messages=[
                    {"role": "system", "content": "You are a business analyst specializing in medical document analysis."},
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"}
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            print(f"Error generating insights: {e}")
            return None

class DataVisualizer:
    def create_condition_distribution(self, df: pd.DataFrame) -> go.Figure:
        # Count conditions by category
        categories = ['initiation_conditions', 'prescribing_conditions', 
                     'pricing_conditions', 'feasibility_conditions']
        condition_counts = []
        
        for category in categories:
            count = df[category].apply(lambda x: len(str(x).split('\n\n'))).sum()
            condition_counts.append(count)
        
        fig = go.Figure(data=[go.Bar(
            x=[cat.split('_')[0].capitalize() for cat in categories],
            y=condition_counts,
            marker_color=['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728']
        )])
        
        fig.update_layout(
            title='Distribution of Conditions by Category',
            xaxis_title='Category',
            yaxis_title='Number of Conditions',
            template='plotly_white'
        )
        return fig

    def create_reason_wordcloud(self, df: pd.DataFrame) -> plt.Figure:
        from wordcloud import WordCloud
        
        # Combine all reasons
        all_reasons = ' '.join([
            str(text) for col in df.columns if 'reasons' in col
            for text in df[col].dropna()
        ])
        
        # Create and generate a word cloud image
        wordcloud = WordCloud(width=800, height=400, background_color='white').generate(all_reasons)
        
        # Create figure
        fig, ax = plt.subplots(figsize=(10, 5))
        ax.imshow(wordcloud, interpolation='bilinear')
        ax.axis('off')
        plt.title('Common Words in Reasons')
        return fig

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        
    def add_title_slide(self, title: str, subtitle: str):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_text_slide(self, title: str, content: List[str]):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        text_frame = slide.placeholders[1].text_frame
        for item in content:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)

    def add_image_slide(self, title: str, image_path: str):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        img_path = image_path
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

    def save_presentation(self) -> BytesIO:
        pptx_buffer = BytesIO()
        self.prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer

def generate_report(data: pd.DataFrame, insights: Dict) -> BytesIO:
    ppt = PPTGenerator()
    
    # Add title slide
    ppt.add_title_slide(
        "Medical Document Analysis Report",
        f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
    )
    
    # Add insights slides
    ppt.add_text_slide("Key Findings", insights['key_findings'])
    ppt.add_text_slide("Business Impact", insights['business_impact'])
    ppt.add_text_slide("Recommendations", insights['recommendations'])
    ppt.add_text_slide("Risk Factors", insights['risk_factors'])
    ppt.add_text_slide("Opportunities", insights['opportunities'])
    
    return ppt.save_presentation()


class ConversationManager:
    def __init__(self):
        self.conversations_dir = "conversations"
        os.makedirs(self.conversations_dir, exist_ok=True)
        
    def save_conversation(self, conversation: List[Dict], name: str = None):
        """Save conversation to file"""
        if name is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            name = f"conversation_{timestamp}"
        
        filepath = os.path.join(self.conversations_dir, f"{name}.json")
        with open(filepath, 'w') as f:
            json.dump({
                'name': name,
                'timestamp': datetime.now().isoformat(),
                'messages': conversation
            }, f)
        return name
    
    def load_conversation(self, name: str) -> List[Dict]:
        """Load conversation from file"""
        filepath = os.path.join(self.conversations_dir, f"{name}.json")
        if os.path.exists(filepath):
            with open(filepath, 'r') as f:
                data = json.load(f)
                return data['messages']
        return []
    
    def list_conversations(self) -> List[str]:
        """List all saved conversations"""
        conversations = []
        for filename in os.listdir(self.conversations_dir):
            if filename.endswith('.json'):
                with open(os.path.join(self.conversations_dir, filename), 'r') as f:
                    data = json.load(f)
                    conversations.append({
                        'name': data['name'],
                        'timestamp': data['timestamp']
                    })
        return sorted(conversations, key=lambda x: x['timestamp'], reverse=True)


class PDFProcessorWithRAG:
    def __init__(self):
        api_key = "sk-proj-2"
        if api_key is None:
            raise ValueError("API key not found in environment variables.")
        self.client = OpenAI(api_key=api_key)
        self.embeddings = OpenAIEmbeddings(openai_api_key=api_key)
        self.keywords = ["Reimbursement", "Recommendation", "Condition", "Reason"]
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self.vector_store = None
        self.processed_data = None
        self.conversations_dir = "conversations"
        os.makedirs(self.conversations_dir, exist_ok=True)
    
    def extract_text_from_pdf(self, pdf_path: str) -> tuple[List[tuple], str]:
        tables_with_keywords = []
        full_text = ""
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                for table in tables:
                    if table and any(any(keyword in str(cell) for cell in table[0]) for keyword in self.keywords):
                        tables_with_keywords.append((page_num, table))
                full_text += page.extract_text() + "\n"
        return tables_with_keywords, full_text

    def create_vector_store(self, texts: List[str]):
        documents = [Document(page_content=text) for text in texts]
        text_chunks = self.text_splitter.split_documents(documents)
        self.vector_store = FAISS.from_documents(text_chunks, self.embeddings)

    def get_relevant_context(self, query: str, k: int = 3) -> str:
        if not self.vector_store:
            return ""
        relevant_docs = self.vector_store.similarity_search(query, k=k)
        return "\n".join(doc.page_content for doc in relevant_docs)

    def format_table_for_llm(self, table_data: List[tuple]) -> str:
        formatted_text = ""
        for page_num, table in table_data:
            formatted_text += f"Page {page_num}:\n"
            for row in table:
                formatted_text += f"{row}\n"
            formatted_text += "\n"
        return formatted_text

    def process_with_llm(self, table_text: str, context: str) -> Dict:
        prompt = f"""
        Using the following context from the medical document and table data, analyze and extract all reimbursement conditions and their reasons exactly as they appear in Table 1.
        Context from document: {context}
        Table data: {table_text}

        Find all conditions listed under each category in Table 1 (Reimbursement Conditions and Reasons) and return them exactly as written.
        For each category (Initiation, Prescribing, Pricing, Feasibility), include all numbered and sub-numbered conditions.

        Return the results in this JSON format:
        {{
            "drug_name": "name of the drug if mentioned",
            "reimbursement_conditions": {{
                "initiation": [
                    {{
                        "condition": "complete condition text including any sub-points",
                        "reason": "exact reason as stated in the table"
                    }}
                ],
                "prescribing": [
                    {{
                        "condition": "complete condition text including any sub-points",
                        "reason": "exact reason as stated in the table"
                    }}
                ],
                "pricing": [
                    {{
                        "condition": "complete condition text including any sub-points",
                        "reason": "exact reason as stated in the table"
                    }}
                ],
                "feasibility": [
                    {{
                        "condition": "complete condition text including any sub-points",
                        "reason": "exact reason as stated in the table"
                    }}
                ]
            }}
        }}

        Important:
        - Capture all conditions exactly as written, including numbering and sub-numbering
        - Include complete condition text for each entry
        - For initiation conditions, include both positive (when treatment should be initiated) and negative (when treatment must not be initiated) conditions
        - Preserve the exact wording from the document
        """
        try:
            response = self.client.chat.completions.create(
                model="gpt-4-turbo-preview",
                messages=[
                    {"role": "system", "content": "You are a medical document analyzer that extracts reimbursement conditions and reasons exactly as they appear in regulatory documents."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0,
                response_format={"type": "json_object"}
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            print(f"Error processing with LLM: {e}")
            return None
        
    def convert_to_dataframe(self, processed_data: Dict, pdf_name: str) -> pd.DataFrame:
        rows = []
        row = {'pdf_name': pdf_name}
        
        # Process each category of conditions
        for category, items in processed_data['reimbursement_conditions'].items():
            conditions = []
            reasons = []
            
            for item in items:
                conditions.append(item['condition'])
                reasons.append(item['reason'])
            
            row[f'{category}_conditions'] = "\n\n".join(conditions)
            row[f'{category}_reasons'] = "\n\n".join(reasons)
        
        rows.append(row)
        return pd.DataFrame(rows)

    def process_multiple_pdfs(self, pdf_path: str):
        all_results = []
        all_texts = []

        try:
            _, full_text = self.extract_text_from_pdf(pdf_path)
            all_texts.append(full_text)
        except Exception as e:
            st.error(f"Error extracting text from {pdf_path}: {e}")
            return None

        self.create_vector_store(all_texts)

        try:
            tables, _ = self.extract_text_from_pdf(pdf_path)
            if not tables:
                st.warning(f"No relevant tables found in {pdf_path}")
                return None

            table_text = self.format_table_for_llm(tables)
            context = self.get_relevant_context(table_text)
            processed_data = self.process_with_llm(table_text, context)
            
            if processed_data:
                df = self.convert_to_dataframe(processed_data, pdf_name=os.path.basename(pdf_path))
                all_results.append(df)
        except Exception as e:
            st.error(f"Error processing {pdf_path}: {e}")
            return None

        if all_results:
            final_df = pd.concat(all_results, ignore_index=True)
            return final_df
        return None

    def process_query(self, query: str, conversation_history: List[Dict]) -> str:
        if self.processed_data is None:
            return "Please upload and process PDF files first."

        # Get relevant context from the processed data
        df_context = self.processed_data.to_string()
        
        conversation_prompt = "\n".join([
            f"{msg['role']}: {msg['content']}" 
            for msg in conversation_history
        ])

        prompt = f"""
        Based on the following processed medical document data and conversation history, 
        please answer the user's question. Provide specific references to the data when possible.

        Processed Data:
        {df_context}

        Conversation History:
        {conversation_prompt}

        User Question: {query}
        """

        try:
            response = self.client.chat.completions.create(
                model="gpt-4-turbo-preview",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant analyzing medical documents. Provide clear, specific answers based on the processed data."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error processing query: {e}"
        
class PDFProcessorWithConversations(PDFProcessorWithRAG):
    def __init__(self):
        super().__init__()
        self.conversation_manager = ConversationManager()


def main():
    st.title("Medical Document Analyzer with Conversation")

    # Initialize session state for conversation history
    if 'conversation_history' not in st.session_state:
        st.session_state.conversation_history = []
    if 'processor' not in st.session_state:
        st.session_state.processor = PDFProcessorWithConversations()
    if 'current_conversation' not in st.session_state:
        st.session_state.current_conversation = None
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = set()
    if 'business_insights' not in st.session_state:
        st.session_state.business_insights = None
    if 'processed_data' not in st.session_state:
        st.session_state.processed_data = None

    with st.sidebar:
        st.header("Navigation")
        if st.button("New Conversation"):
            st.session_state.conversation_history = []
            st.session_state.current_conversation = None
            st.rerun()

    # Main content area with tabs
    tab1, tab2, tab3, tab4 = st.tabs(["Document Processing", "Visualizations", "Business Insights", "Report Generation"])
    
    # Document Processing Tab
    with tab1:
        st.header("Upload and Process Documents")
        uploaded_files = st.file_uploader("Upload PDF Files", accept_multiple_files=True, type='pdf')
        
        if uploaded_files:
            if st.button("Process PDFs"):    
                with st.spinner("Processing PDFs..."):
                    try:
                        all_results = []
                        for pdf_file in uploaded_files:
                            temp_pdf_path = f"temp_{pdf_file.name}"
                            with open(temp_pdf_path, "wb") as f:
                                f.write(pdf_file.read())
                            
                            processed_df = st.session_state.processor.process_multiple_pdfs(temp_pdf_path)
                            if processed_df is not None:
                                all_results.append(processed_df)
                            os.remove(temp_pdf_path)
                        
                        if all_results:
                            final_df = pd.concat(all_results, ignore_index=True)
                            st.session_state.processed_data = final_df
                            st.success("PDFs processed successfully!")
                            
                            # Display the processed data
                            st.dataframe(final_df)
                            
                            # Create download button for CSV
                            csv = final_df.to_csv(index=False)
                            st.download_button(
                                label="Download Processed Data as CSV",
                                data=csv,
                                file_name="processed_data.csv",
                                mime="text/csv"
                            )
                    except Exception as e:
                        st.error(f"Error processing PDFs: {e}")
    
    # Visualizations Tab
    with tab2:
        st.header("Data Visualizations")
        if st.session_state.processed_data is not None:
            visualizer = DataVisualizer()
            
            st.subheader("Condition Distribution")
            try:
                fig1 = visualizer.create_condition_distribution(st.session_state.processed_data)
                st.plotly_chart(fig1, use_container_width=True)
            except Exception as e:
                st.error(f"Error creating condition distribution: {e}")
            
            st.subheader("Reasons Word Cloud")
            try:
                fig2 = visualizer.create_reason_wordcloud(st.session_state.processed_data)
                st.pyplot(fig2)
            except Exception as e:
                st.error(f"Error creating word cloud: {e}")
        else:
            st.info("Please process documents first to view visualizations.")
    
    # Business Insights Tab
    with tab3:
        st.header("Business Insights")
        if st.session_state.processed_data is not None:
            if st.button("Generate Business Insights"):
                with st.spinner("Generating insights..."):
                    try:
                        business_analyzer = BusinessAnalyzer(st.session_state.processor.client)
                        insights = business_analyzer.generate_insights(st.session_state.processed_data)
                        st.session_state.business_insights = insights
                        
                        if insights:
                            for category in ['key_findings', 'business_impact', 'recommendations', 'risk_factors', 'opportunities']:
                                st.subheader(category.replace('_', ' ').title())
                                for item in insights[category]:
                                    st.write(f"• {item}")
                    except Exception as e:
                        st.error(f"Error generating insights: {e}")
        else:
            st.info("Please process documents first to generate business insights.")
    
    # Report Generation Tab
    with tab4:
        st.header("Report Generation")
        if st.session_state.processed_data is not None and st.session_state.business_insights is not None:
            if st.button("Generate PowerPoint Report"):
                with st.spinner("Generating report..."):
                    try:
                        pptx_buffer = generate_report(
                            st.session_state.processed_data,
                            st.session_state.business_insights
                        )
                        
                        st.download_button(
                            label="Download PowerPoint Report",
                            data=pptx_buffer,
                            file_name="medical_document_analysis.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Error generating report: {e}")
        else:
            st.info("Please process documents and generate business insights first to create a report.")

    # Chat Interface
    st.header("Chat Interface")
    if st.session_state.current_conversation:
        st.subheader(f"Current Conversation: {st.session_state.current_conversation}")

    # Display conversation history
    for message in st.session_state.conversation_history:
        with st.chat_message(message["role"]):
            st.write(message["content"])

    # Conversation name input
    conversation_name = st.text_input(
        "Conversation Name (optional)",
        key="conversation_name",
        placeholder="Enter a name to save this conversation"
    )

    # Query input
    if query := st.chat_input("Ask a question about the processed documents"):
        if not st.session_state.processed_data is not None:
            st.warning("Please process some PDFs first before asking questions.")
            return

        st.session_state.conversation_history.append({"role": "user", "content": query})
        
        with st.spinner("Generating response..."):
            try:
                response = st.session_state.processor.process_query(
                    query, 
                    st.session_state.conversation_history
                )
                
                st.session_state.conversation_history.append({"role": "assistant", "content": response})
                
                name = conversation_name if conversation_name else None
                st.session_state.current_conversation = (
                    st.session_state.processor.conversation_manager.save_conversation(
                        st.session_state.conversation_history,
                        name=name
                    )
                )
            except Exception as e:
                st.error(f"Error processing query: {e}")
        
        st.rerun()

if __name__ == "__main__":
    main()